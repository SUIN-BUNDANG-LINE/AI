import tempfile
from http import HTTPStatus

import requests
from langchain.docstore.document import Document
from langchain_community.document_loaders import (
    PyMuPDFLoader,
    Docx2txtLoader,
)
from pptx import Presentation

from app.core.error.business_exception import business_exception
from app.core.error.error_code import ErrorCode
from app.core.util.file_manager import FileManager


class DocumentManager:
    def __init__(self):
        self.pdf_loader = PyMuPDFLoader
        self.docx_loader = Docx2txtLoader

    def text_from_file_url(self, file_url: str):
        extension = FileManager.get_file_extension_from_url(file_url)
        match extension:
            case ".pdf":
                return self.text_from_pdf_file_url(file_url)
            case ".txt":
                return self.text_from_txt_file_url(file_url)
            case ".docx":
                return self.text_from_docx_file_url(file_url)
            case ".pptx":
                return self.text_from_pptx_file_url(file_url)
            case _:
                raise business_exception(ErrorCode.FILE_EXTENSION_NOT_SUPPORTED)

    def text_from_pdf_file_url(self, file_url: str):
        documents = self.pdf_loader(file_url).load()
        return DocumentManager.__documents_to_text(documents)

    def text_from_docx_file_url(self, file_url: str):
        documents = self.docx_loader(file_url).load()
        return DocumentManager.__documents_to_text(documents)

    @staticmethod
    def text_from_pptx_file_url(file_url: str):
        response = requests.get(file_url)
        response.raise_for_status()

        # 임시 파일 생성
        with tempfile.NamedTemporaryFile(delete=True, suffix=".pptx") as temp_file:
            temp_file.write(response.content)
            temp_file.flush()

            presentation = Presentation(temp_file.name)
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

            return "\n".join(text)

    @staticmethod
    def text_from_txt_file_url(file_url: str):
        response = requests.get(file_url)
        if response.status_code != HTTPStatus.OK:
            raise business_exception(ErrorCode.FILE_NOT_FOUND)
        response.encoding = "utf-8"
        text_content = response.text

        documents = [Document(page_content=text_content)]
        return DocumentManager.__documents_to_text(documents)

    @staticmethod
    def __documents_to_text(documents: list[Document]):
        documents_text = ""
        for document in documents:
            documents_text += f"""
            ----------------------------
            {document.page_content}
            """

        return documents_text
